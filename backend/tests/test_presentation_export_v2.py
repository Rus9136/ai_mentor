"""
Tests for presentation_export_v2.py — smoke tests per slide type,
backward compatibility, edge cases.

Runs WITHOUT database — pure PPTX generation.
"""
import json
from io import BytesIO
from pathlib import Path

import pytest
from pptx import Presentation

from app.services.presentation_export_v2 import (
    THEMES,
    THEME_ALIASES,
    ImageProvider,
    export_to_pptx,
    get_available_templates,
    _resolve_theme,
)


SAMPLE_SLIDES_PATH = Path(__file__).parent.parent.parent / "sample_slides.json"

CONTEXT = {
    "paragraph_title": "Қасым хан",
    "chapter_title": "XVI ғасыр",
    "textbook_title": "Қазақстан тарихы 7",
    "subject": "Қазақстан тарихы",
    "grade_level": 7,
    "textbook_id": 25,
    "theme": "warm",
}

# ═══════════════════════════════════════════════════════════
#  THEME RESOLUTION
# ═══════════════════════════════════════════════════════════


class TestThemeResolution:
    def test_resolve_from_context(self):
        theme = _resolve_theme({}, {"theme": "midnight"}, None)
        assert theme.name == "midnight"

    def test_resolve_from_template_param(self):
        theme = _resolve_theme({}, {}, "forest")
        assert theme.name == "forest"

    def test_alias_green_to_forest(self):
        theme = _resolve_theme({}, {}, "green")
        assert theme.name == "forest"

    def test_alias_blue_to_warm(self):
        theme = _resolve_theme({}, {}, "blue")
        assert theme.name == "warm"

    def test_unknown_theme_defaults_warm(self):
        theme = _resolve_theme({}, {}, "nonexistent")
        assert theme.name == "warm"

    def test_context_overrides_template(self):
        theme = _resolve_theme({}, {"theme": "midnight"}, "forest")
        assert theme.name == "midnight"


# ═══════════════════════════════════════════════════════════
#  SMOKE TESTS — each slide type renders without error
# ═══════════════════════════════════════════════════════════


def _make_slides_data(*slides):
    return {"title": "Test", "slides": list(slides)}


class TestSlideSmoke:
    def test_title_slide(self):
        data = _make_slides_data({"type": "title", "title": "Тест", "subtitle": "Ішкі тест"})
        buf = export_to_pptx(data, CONTEXT)
        prs = Presentation(buf)
        assert len(prs.slides) == 1

    def test_objectives_slide(self):
        data = _make_slides_data({"type": "objectives", "title": "Мақсат", "items": ["A", "B", "C"]})
        buf = export_to_pptx(data, CONTEXT)
        assert len(Presentation(buf).slides) == 1

    def test_content_image_left(self):
        data = _make_slides_data({
            "type": "content", "title": "Topic", "body": "Body text",
            "layout_hint": "image_left", "image_query": "kazakh steppe",
        })
        buf = export_to_pptx(data, CONTEXT)
        assert len(Presentation(buf).slides) == 1

    def test_content_image_right(self):
        data = _make_slides_data({
            "type": "content", "title": "Topic", "body": "Body text",
            "layout_hint": "image_right",
        })
        buf = export_to_pptx(data, CONTEXT)
        assert len(Presentation(buf).slides) == 1

    def test_content_stat_callout(self):
        data = _make_slides_data({
            "type": "content", "title": "Key Fact", "body": "Body text",
            "layout_hint": "stat_callout", "stat_value": "1511",
            "stat_label": "year of coronation",
        })
        buf = export_to_pptx(data, CONTEXT)
        assert len(Presentation(buf).slides) == 1

    def test_content_no_layout_hint(self):
        """V1-style content without layout_hint — should default to image_left."""
        data = _make_slides_data({"type": "content", "title": "Old", "body": "Old body"})
        buf = export_to_pptx(data, CONTEXT)
        assert len(Presentation(buf).slides) == 1

    def test_key_terms_slide(self):
        data = _make_slides_data({
            "type": "key_terms", "title": "Terms",
            "terms": [
                {"term": "A", "definition": "Def A"},
                {"term": "B", "definition": "Def B"},
                {"term": "C", "definition": "Def C"},
                {"term": "D", "definition": "Def D"},
            ],
        })
        buf = export_to_pptx(data, CONTEXT)
        assert len(Presentation(buf).slides) == 1

    def test_quiz_slide(self):
        data = _make_slides_data({
            "type": "quiz", "question": "What?",
            "options": ["A", "B", "C", "D"], "answer": 2,
        })
        buf = export_to_pptx(data, CONTEXT)
        prs = Presentation(buf)
        assert len(prs.slides) == 1
        # Check answer is in speaker notes
        notes = prs.slides[0].notes_slide.notes_text_frame.text
        assert "C" in notes

    def test_summary_slide(self):
        data = _make_slides_data({
            "type": "summary", "title": "Conclusion",
            "items": ["Point 1", "Point 2", "Point 3"],
        })
        buf = export_to_pptx(data, CONTEXT)
        assert len(Presentation(buf).slides) == 1

    def test_unknown_type_falls_back(self):
        data = _make_slides_data({"type": "unknown_type", "title": "?", "body": "fallback"})
        buf = export_to_pptx(data, CONTEXT)
        assert len(Presentation(buf).slides) == 1


# ═══════════════════════════════════════════════════════════
#  ALL THEMES RENDER
# ═══════════════════════════════════════════════════════════


class TestAllThemes:
    @pytest.mark.parametrize("theme_name", list(THEMES.keys()))
    def test_full_deck_each_theme(self, theme_name):
        data = _make_slides_data(
            {"type": "title", "title": "T", "subtitle": "S"},
            {"type": "objectives", "title": "O", "items": ["A"]},
            {"type": "content", "title": "C", "body": "B", "layout_hint": "image_left"},
            {"type": "content", "title": "C2", "body": "B2", "layout_hint": "stat_callout",
             "stat_value": "1511", "stat_label": "year"},
            {"type": "key_terms", "title": "KT", "terms": [{"term": "X", "definition": "Y"}]},
            {"type": "quiz", "question": "Q?", "options": ["A", "B", "C", "D"], "answer": 0},
            {"type": "summary", "title": "S", "items": ["P1"]},
        )
        ctx = {**CONTEXT, "theme": theme_name}
        buf = export_to_pptx(data, ctx)
        prs = Presentation(buf)
        assert len(prs.slides) == 7


# ═══════════════════════════════════════════════════════════
#  BACKWARD COMPATIBILITY — V1 JSON
# ═══════════════════════════════════════════════════════════


class TestBackwardCompat:
    def test_v1_json_renders(self):
        """Old v1 JSON without layout_hint/image_query/stat fields."""
        v1_data = {
            "title": "Test v1",
            "slides": [
                {"type": "title", "title": "V1 Title", "subtitle": "V1 Subtitle"},
                {"type": "objectives", "title": "Goals", "items": ["Goal 1", "Goal 2"]},
                {"type": "content", "title": "Content", "body": "Some text."},
                {"type": "content", "title": "Content2", "body": "More text.",
                 "image_url": "/uploads/textbook-images/25/nonexistent.jpg"},
                {"type": "key_terms", "title": "Terms", "terms": [
                    {"term": "T1", "definition": "D1"},
                    {"term": "T2", "definition": "D2"},
                ]},
                {"type": "quiz", "title": "Quiz", "question": "Q?",
                 "options": ["A", "B", "C", "D"], "answer": 1},
                {"type": "summary", "title": "Sum", "items": ["I1", "I2"]},
            ],
        }
        buf = export_to_pptx(v1_data, CONTEXT)
        prs = Presentation(buf)
        assert len(prs.slides) == 7


# ═══════════════════════════════════════════════════════════
#  EDGE CASES
# ═══════════════════════════════════════════════════════════


class TestEdgeCases:
    def test_empty_slides(self):
        buf = export_to_pptx({"title": "Empty", "slides": []}, CONTEXT)
        prs = Presentation(buf)
        assert len(prs.slides) == 0

    def test_stat_value_long(self):
        """Long stat_value should render without overflow."""
        data = _make_slides_data({
            "type": "content", "title": "Long", "body": "B",
            "layout_hint": "stat_callout",
            "stat_value": "1234567",
            "stat_label": "very long label text here",
        })
        buf = export_to_pptx(data, CONTEXT)
        assert len(Presentation(buf).slides) == 1

    def test_cyrillic_diacritics(self):
        """Kazakh diacritics: Қ, Ұ, Ң, Ғ, Ү, Ө, І, Ә."""
        data = _make_slides_data({
            "type": "content", "title": "Қасым Ұлытау Ңәзір",
            "body": "Ғылым Үкімет Өнер Іс Әділет",
            "layout_hint": "image_left",
        })
        buf = export_to_pptx(data, CONTEXT)
        assert len(Presentation(buf).slides) == 1

    def test_missing_context(self):
        """None context should not crash."""
        data = _make_slides_data({"type": "title", "title": "T", "subtitle": "S"})
        buf = export_to_pptx(data, {})
        assert len(Presentation(buf).slides) == 1

    def test_quiz_invalid_answer_index(self):
        """answer index beyond options length should not crash."""
        data = _make_slides_data({
            "type": "quiz", "question": "Q?",
            "options": ["A", "B"], "answer": 99,
        })
        buf = export_to_pptx(data, CONTEXT)
        assert len(Presentation(buf).slides) == 1

    def test_key_terms_empty_list(self):
        data = _make_slides_data({"type": "key_terms", "title": "T", "terms": []})
        buf = export_to_pptx(data, CONTEXT)
        assert len(Presentation(buf).slides) == 1


# ═══════════════════════════════════════════════════════════
#  SAMPLE_SLIDES.JSON INTEGRATION
# ═══════════════════════════════════════════════════════════


class TestSampleSlides:
    def test_sample_json(self):
        if not SAMPLE_SLIDES_PATH.exists():
            pytest.skip("sample_slides.json not found")
        with open(SAMPLE_SLIDES_PATH, encoding="utf-8") as f:
            data = json.load(f)
        buf = export_to_pptx(data, CONTEXT)
        prs = Presentation(buf)
        assert len(prs.slides) == 10
        # All slides should have shapes (no empty slides)
        for i, slide in enumerate(prs.slides):
            assert len(slide.shapes) > 0, f"Slide {i+1} has no shapes"


# ═══════════════════════════════════════════════════════════
#  IMAGE PROVIDER
# ═══════════════════════════════════════════════════════════


class TestImageProvider:
    def test_no_key_returns_none(self):
        provider = ImageProvider(unsplash_key="")
        assert provider.fetch("kazakh steppe") is None

    def test_empty_query_returns_none(self):
        provider = ImageProvider(unsplash_key="fake")
        assert provider.fetch("") is None

    def test_cache_stores_results(self):
        provider = ImageProvider(unsplash_key="")
        # Class-level cache — inject directly
        ImageProvider._cache["test_cached"] = (b"fake_bytes", 9999999999)
        assert provider.fetch("test_cached") == b"fake_bytes"
        # Cleanup
        del ImageProvider._cache["test_cached"]


# ═══════════════════════════════════════════════════════════
#  TEMPLATES API
# ═══════════════════════════════════════════════════════════


class TestTemplatesAPI:
    def test_available_templates(self):
        templates = get_available_templates()
        slugs = [t["slug"] for t in templates]
        assert "warm" in slugs
        assert "forest" in slugs
        assert "midnight" in slugs

    def test_aliases(self):
        assert THEME_ALIASES["green"] == "forest"
        assert THEME_ALIASES["blue"] == "warm"
