"""
PPTX export v2 — Gamma-style exporter for AI-Mentor.

Key improvements over v1:
- Full theme palettes with motif system (not just accent color swaps)
- 3 layout variants per content slide (image_left, image_right, stat_callout)
- Image integration via Unsplash with graceful fallback to colored placeholders
- Redesigned key_terms (card grid), quiz (2x2), summary (dark sandwich)
- Speaker notes on every slide
- Backward compatible: old v1 JSON renders fine (missing fields get defaults)

Public API matches v1 signature:
    export_to_pptx(slides_data, context_data, template) -> BytesIO
"""
from __future__ import annotations

import io
import logging
import time
from dataclasses import dataclass
import httpx
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Emu, Inches, Pt

logger = logging.getLogger(__name__)


# ─────────────────────────────────────────────────────────────────────────────
# THEMES — full palettes with dominance + motif
# ─────────────────────────────────────────────────────────────────────────────


@dataclass
class Theme:
    name: str
    bg: RGBColor
    surface: RGBColor
    primary: RGBColor
    primary_dark: RGBColor
    accent: RGBColor
    text: RGBColor
    text_muted: RGBColor
    text_on_primary: RGBColor
    motif: str  # 'rounded_card' | 'side_bar' | 'circle_badge'
    header_font: str = "Calibri"
    body_font: str = "Calibri"
    label_ru: str = ""
    label_kk: str = ""


FONT_FAMILY = "Calibri"  # legacy default, prefer theme.body_font

THEMES: dict[str, Theme] = {
    "warm": Theme(
        name="warm",
        bg=RGBColor(0xFF, 0xF8, 0xF2),
        surface=RGBColor(0xFF, 0xFF, 0xFF),
        primary=RGBColor(0xE8, 0x81, 0x34),
        primary_dark=RGBColor(0xC2, 0x65, 0x1F),
        accent=RGBColor(0x9F, 0xC9, 0x69),
        text=RGBColor(0x2B, 0x1D, 0x12),
        text_muted=RGBColor(0x8A, 0x73, 0x60),
        text_on_primary=RGBColor(0xFF, 0xFF, 0xFF),
        motif="rounded_card",
        label_ru="Тёплый",
        label_kk="Жылы",
    ),
    "forest": Theme(
        name="forest",
        bg=RGBColor(0xF4, 0xF8, 0xF1),
        surface=RGBColor(0xFF, 0xFF, 0xFF),
        primary=RGBColor(0x2C, 0x5F, 0x2D),
        primary_dark=RGBColor(0x1A, 0x40, 0x1B),
        accent=RGBColor(0xF5, 0xC2, 0x3E),
        text=RGBColor(0x16, 0x2A, 0x18),
        text_muted=RGBColor(0x5E, 0x73, 0x5F),
        text_on_primary=RGBColor(0xFF, 0xFF, 0xFF),
        motif="side_bar",
        label_ru="Изумрудный",
        label_kk="Зүмірет",
    ),
    "midnight": Theme(
        name="midnight",
        bg=RGBColor(0x10, 0x14, 0x28),
        surface=RGBColor(0x1B, 0x21, 0x3D),
        primary=RGBColor(0x6E, 0x8B, 0xFF),
        primary_dark=RGBColor(0x4F, 0x67, 0xCC),
        accent=RGBColor(0xFF, 0xC8, 0x4C),
        text=RGBColor(0xEC, 0xEE, 0xF7),
        text_muted=RGBColor(0x9A, 0xA1, 0xC0),
        text_on_primary=RGBColor(0x10, 0x14, 0x28),
        motif="circle_badge",
        label_ru="Полночь",
        label_kk="Түн",
    ),
    "parchment": Theme(
        name="parchment",
        bg=RGBColor(0xFA, 0xF3, 0xE8),
        surface=RGBColor(0xFF, 0xFF, 0xFF),
        primary=RGBColor(0x6D, 0x2E, 0x46),
        primary_dark=RGBColor(0x4E, 0x1F, 0x32),
        accent=RGBColor(0xD4, 0xA5, 0x74),
        text=RGBColor(0x2C, 0x18, 0x10),
        text_muted=RGBColor(0x7A, 0x65, 0x58),
        text_on_primary=RGBColor(0xFF, 0xFF, 0xFF),
        motif="side_bar",
        header_font="Georgia",
        label_ru="Пергамент",
        label_kk="Пергамент",
    ),
    "slate": Theme(
        name="slate",
        bg=RGBColor(0xF5, 0xF5, 0xF7),
        surface=RGBColor(0xFF, 0xFF, 0xFF),
        primary=RGBColor(0x36, 0x45, 0x4F),
        primary_dark=RGBColor(0x24, 0x2F, 0x36),
        accent=RGBColor(0x6B, 0x8F, 0x71),
        text=RGBColor(0x1D, 0x1D, 0x1F),
        text_muted=RGBColor(0x86, 0x86, 0x8B),
        text_on_primary=RGBColor(0xFF, 0xFF, 0xFF),
        motif="rounded_card",
        label_ru="Минимализм",
        label_kk="Минимализм",
    ),
    "electric": Theme(
        name="electric",
        bg=RGBColor(0xF8, 0xFA, 0xFC),
        surface=RGBColor(0xFF, 0xFF, 0xFF),
        primary=RGBColor(0x25, 0x63, 0xEB),
        primary_dark=RGBColor(0x1D, 0x4E, 0xD8),
        accent=RGBColor(0xF5, 0x9E, 0x0B),
        text=RGBColor(0x0F, 0x17, 0x2A),
        text_muted=RGBColor(0x64, 0x74, 0x8B),
        text_on_primary=RGBColor(0xFF, 0xFF, 0xFF),
        motif="circle_badge",
        label_ru="Электрик",
        label_kk="Электрик",
    ),
    "lavender": Theme(
        name="lavender",
        bg=RGBColor(0xF5, 0xF0, 0xFF),
        surface=RGBColor(0xFF, 0xFF, 0xFF),
        primary=RGBColor(0x7C, 0x3A, 0xED),
        primary_dark=RGBColor(0x6D, 0x28, 0xD9),
        accent=RGBColor(0xF4, 0x72, 0xB6),
        text=RGBColor(0x1E, 0x1B, 0x4B),
        text_muted=RGBColor(0x6B, 0x72, 0x80),
        text_on_primary=RGBColor(0xFF, 0xFF, 0xFF),
        motif="rounded_card",
        label_ru="Лаванда",
        label_kk="Лаванда",
    ),
    "coral": Theme(
        name="coral",
        bg=RGBColor(0xFF, 0xF5, 0xF2),
        surface=RGBColor(0xFF, 0xFF, 0xFF),
        primary=RGBColor(0xE8, 0x55, 0x3D),
        primary_dark=RGBColor(0xC2, 0x41, 0x2D),
        accent=RGBColor(0xFB, 0xBF, 0x24),
        text=RGBColor(0x27, 0x15, 0x0E),
        text_muted=RGBColor(0x9A, 0x7B, 0x70),
        text_on_primary=RGBColor(0xFF, 0xFF, 0xFF),
        motif="rounded_card",
        label_ru="Коралл",
        label_kk="Коралл",
    ),
    "ocean": Theme(
        name="ocean",
        bg=RGBColor(0xEF, 0xF6, 0xFF),
        surface=RGBColor(0xFF, 0xFF, 0xFF),
        primary=RGBColor(0x06, 0x5A, 0x82),
        primary_dark=RGBColor(0x04, 0x45, 0x66),
        accent=RGBColor(0x0E, 0xA5, 0xE9),
        text=RGBColor(0x0C, 0x23, 0x40),
        text_muted=RGBColor(0x5B, 0x7A, 0x99),
        text_on_primary=RGBColor(0xFF, 0xFF, 0xFF),
        motif="side_bar",
        header_font="Georgia",
        label_ru="Океан",
        label_kk="Мұхит",
    ),
    "sage": Theme(
        name="sage",
        bg=RGBColor(0xF2, 0xF7, 0xF2),
        surface=RGBColor(0xFF, 0xFF, 0xFF),
        primary=RGBColor(0x5F, 0x71, 0x61),
        primary_dark=RGBColor(0x4A, 0x5A, 0x4C),
        accent=RGBColor(0xD4, 0xA5, 0x74),
        text=RGBColor(0x1A, 0x2E, 0x1A),
        text_muted=RGBColor(0x6B, 0x7C, 0x6B),
        text_on_primary=RGBColor(0xFF, 0xFF, 0xFF),
        motif="rounded_card",
        label_ru="Шалфей",
        label_kk="Жусан",
    ),
}

# Alias: v1 "green" maps to v2 "forest"
THEME_ALIASES: dict[str, str] = {
    "green": "forest",
    "blue": "warm",  # v1 default → warm
}


# ─────────────────────────────────────────────────────────────────────────────
# LOCALIZED UI STRINGS — eyebrow text per language
# ─────────────────────────────────────────────────────────────────────────────

_L10N: dict[str, dict[str, str]] = {
    "kk": {
        "lesson": "САБАҚ",
        "topic": "ТАҚЫРЫП",
        "key_fact": "НЕГІЗГІ ФАКТ",
        "terms": "ТЕРМИНДЕР",
        "terms_title": "Негізгі ұғымдар",
        "quiz": "БІЛІМДІ ТЕКСЕР",
        "summary": "ҚОРЫТЫНДЫ",
        "summary_title": "Не үйрендік",
        "objectives_title": "Сабақтың мақсаты",
        "class_suffix": "сынып",
        "notes_intro": "Кіріспе слайд",
        "notes_objectives": "Оқушыларға сабақтың мақсаттарын түсіндіріңіз.",
        "notes_discuss": "Талқылау",
        "notes_fact": "Маңызды факт",
        "notes_terms": "Әр терминді 20-30 секундта талқылаңыз.",
        "notes_quiz": "Оқушыларға 30 секунд ойлануға уақыт беріңіз.",
        "notes_quiz_answer": "Дұрыс жауап",
        "notes_summary": "Сабақты қорытындылау.",
    },
    "ru": {
        "lesson": "УРОК",
        "topic": "ТЕМА",
        "key_fact": "КЛЮЧЕВОЙ ФАКТ",
        "terms": "ТЕРМИНЫ",
        "terms_title": "Основные понятия",
        "quiz": "ПРОВЕРЬ ЗНАНИЯ",
        "summary": "ИТОГИ",
        "summary_title": "Что мы узнали",
        "objectives_title": "Цели урока",
        "class_suffix": "класс",
        "notes_intro": "Вводный слайд",
        "notes_objectives": "Объясните ученикам цели урока.",
        "notes_discuss": "Обсуждение",
        "notes_fact": "Важный факт",
        "notes_terms": "Обсудите каждый термин за 20-30 секунд.",
        "notes_quiz": "Дайте ученикам 30 секунд на размышление.",
        "notes_quiz_answer": "Правильный ответ",
        "notes_summary": "Подведение итогов урока.",
    },
}


def _get_l10n(context_data: dict) -> dict[str, str]:
    """Get localized strings from context language, default to kk."""
    lang = (context_data or {}).get("language", "kk")
    return _L10N.get(lang, _L10N["kk"])


# ─────────────────────────────────────────────────────────────────────────────
# SLIDE GEOMETRY — 16:9, generous margins
# ─────────────────────────────────────────────────────────────────────────────

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)
MARGIN = Inches(0.6)
GUTTER = Inches(0.4)


# ─────────────────────────────────────────────────────────────────────────────
# IMAGE PROVIDER — Unsplash with in-memory cache and graceful fallback
# ─────────────────────────────────────────────────────────────────────────────


class ImageProvider:
    """Fetches images by keyword from Unsplash. Falls back to placeholder if no API key."""

    _warned_no_key: bool = False
    _cache: dict[str, tuple[bytes | None, float]] = {}  # class-level — survives across calls
    _cache_ttl: int = 30 * 24 * 3600  # 30 days

    def __init__(self, unsplash_key: str | None = None):
        if unsplash_key is None:
            from app.core.config import settings
            unsplash_key = settings.UNSPLASH_ACCESS_KEY or ""
        self.key = unsplash_key
        self._client: httpx.Client | None = None
        if not self.key and not ImageProvider._warned_no_key:
            logger.warning(
                "UNSPLASH_ACCESS_KEY not set — images will use colored placeholders"
            )
            ImageProvider._warned_no_key = True

    def _get_client(self) -> httpx.Client:
        if self._client is None:
            self._client = httpx.Client(timeout=8.0)
        return self._client

    def fetch(self, query: str) -> bytes | None:
        if not query:
            return None

        # Check class-level cache
        now = time.time()
        if query in ImageProvider._cache:
            data, ts = ImageProvider._cache[query]
            if now - ts < ImageProvider._cache_ttl:
                return data

        if not self.key:
            return None

        try:
            client = self._get_client()
            resp = client.get(
                "https://api.unsplash.com/search/photos",
                params={"query": query, "per_page": 1, "orientation": "landscape"},
                headers={"Authorization": f"Client-ID {self.key}"},
            )
            if resp.status_code in (429, 401, 403):
                logger.warning("Unsplash returned %d, falling back to placeholder", resp.status_code)
                ImageProvider._cache[query] = (None, now)
                return None
            resp.raise_for_status()
            results = resp.json().get("results", [])
            if not results:
                ImageProvider._cache[query] = (None, now)
                return None
            url = results[0]["urls"]["regular"]
            img_resp = client.get(url)
            img_resp.raise_for_status()
            img_bytes = img_resp.content
            ImageProvider._cache[query] = (img_bytes, now)
            return img_bytes
        except Exception:
            logger.warning("Unsplash fetch failed for query=%s", query, exc_info=True)
            return None


# ─────────────────────────────────────────────────────────────────────────────
# DRAWING PRIMITIVES
# ─────────────────────────────────────────────────────────────────────────────


def _rect(slide, x, y, w, h, fill: RGBColor, *, rounded: bool = False,
          corner_radius: float = 0.08, line: RGBColor | None = None):
    shape_type = MSO_SHAPE.ROUNDED_RECTANGLE if rounded else MSO_SHAPE.RECTANGLE
    shape = slide.shapes.add_shape(shape_type, int(x), int(y), int(w), int(h))
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    if line is None:
        shape.line.fill.background()
    else:
        shape.line.color.rgb = line
        shape.line.width = Pt(0.75)
    if rounded:
        try:
            shape.adjustments[0] = corner_radius
        except Exception:
            pass
    shape.shadow.inherit = False
    return shape


def _circle(slide, cx, cy, diameter, fill: RGBColor):
    x = int(cx - diameter / 2)
    y = int(cy - diameter / 2)
    shape = slide.shapes.add_shape(MSO_SHAPE.OVAL, x, y, int(diameter), int(diameter))
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    shape.line.fill.background()
    shape.shadow.inherit = False
    return shape


def _text(slide, x, y, w, h, text: str, *, size: int, color: RGBColor,
          bold: bool = False, align=PP_ALIGN.LEFT,
          anchor=MSO_ANCHOR.TOP, line_spacing: float = 1.15,
          font: str = FONT_FAMILY):
    box = slide.shapes.add_textbox(int(x), int(y), int(w), int(h))
    tf = box.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = Emu(0)
    tf.margin_top = tf.margin_bottom = Emu(0)
    tf.vertical_anchor = anchor
    tf.auto_size = None

    p = tf.paragraphs[0]
    p.alignment = align
    p.line_spacing = line_spacing
    run = p.add_run()
    run.text = text
    run.font.name = font
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    return box


def _badge_text(slide, cx, cy, text: str, *, size: int, color: RGBColor,
                bold: bool = True, font: str = FONT_FAMILY):
    """Centered text on a point — for letters/numbers in circles."""
    cap_h = Pt(size).emu * 0.85
    box_w = Pt(size).emu * max(len(text), 1) * 0.7
    box = slide.shapes.add_textbox(
        int(cx - box_w / 2), int(cy - cap_h * 0.65),
        int(box_w), int(cap_h * 1.3),
    )
    tf = box.text_frame
    tf.word_wrap = False
    tf.margin_left = tf.margin_right = Emu(0)
    tf.margin_top = tf.margin_bottom = Emu(0)
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    tf.auto_size = None
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    p.line_spacing = 1.0
    run = p.add_run()
    run.text = text
    run.font.name = font
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    return box


def _image_fill(slide, image_bytes: bytes, x, y, w, h):
    stream = io.BytesIO(image_bytes)
    return slide.shapes.add_picture(stream, int(x), int(y), width=int(w), height=int(h))


def _image_placeholder(slide, x, y, w, h, theme: Theme):
    """Colored block with decorative circles when no image available."""
    _rect(slide, x, y, w, h, theme.primary, rounded=True, corner_radius=0.04)
    _circle(slide, x + w * 0.7, y + h * 0.3, Inches(2.0), theme.primary_dark)
    _circle(slide, x + w * 0.3, y + h * 0.7, Inches(1.4), theme.accent)


def _speaker_notes(slide, text: str):
    notes = slide.notes_slide.notes_text_frame
    notes.text = text


# ─────────────────────────────────────────────────────────────────────────────
# SLIDE BUILDERS
# ─────────────────────────────────────────────────────────────────────────────


class SlideBuilder:
    def __init__(self, prs: Presentation, theme: Theme, images: ImageProvider,
                 l10n: dict[str, str] | None = None):
        self.prs = prs
        self.t = theme
        self.images = images
        self.l = l10n or _L10N["kk"]

    def _blank(self):
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])
        bg = _rect(slide, 0, 0, SLIDE_W, SLIDE_H, self.t.bg)
        bg.shadow.inherit = False
        return slide

    def _motif(self, slide):
        t = self.t
        if t.motif == "side_bar":
            _rect(slide, 0, 0, Inches(0.18), SLIDE_H, t.primary)
        elif t.motif == "circle_badge":
            _circle(slide, SLIDE_W - Inches(0.6), Inches(0.6), Inches(0.5), t.accent)
        elif t.motif == "rounded_card":
            _rect(slide, SLIDE_W - Inches(0.9), SLIDE_H - Inches(0.3),
                  Inches(0.9), Inches(0.08), t.primary)

    def _footer(self, slide, page_num: int, total: int, deck_title: str):
        _text(
            slide, MARGIN, SLIDE_H - Inches(0.45),
            SLIDE_W - 2 * MARGIN, Inches(0.3),
            f"{deck_title}   \u00b7   {page_num} / {total}",
            size=10, color=self.t.text_muted,
        )

    # ── TITLE ────────────────────────────────────────────────────────────
    def build_title(self, data: dict, context: dict, deck_title: str):
        slide = self._blank()
        t = self.t

        # Half-bleed image on the right
        img_x = SLIDE_W * 0.55
        img_w = SLIDE_W - img_x
        img_bytes = self.images.fetch(data.get("image_query", ""))
        if img_bytes:
            _image_fill(slide, img_bytes, img_x, 0, img_w, SLIDE_H)
        else:
            _image_placeholder(slide, img_x, 0, img_w, SLIDE_H, t)

        # Blending strip
        _rect(slide, img_x, 0, Inches(0.3), SLIDE_H, t.bg)

        left_w = img_x - MARGIN * 2

        # Eyebrow
        _text(slide, MARGIN, Inches(1.4), left_w, Inches(0.4),
              self.l["lesson"], size=14, color=t.primary, bold=True)

        # Title
        _text(slide, MARGIN, Inches(1.9), left_w, Inches(2.6),
              data.get("title", ""), size=54, color=t.text, bold=True,
              line_spacing=1.05)

        # Subtitle
        if data.get("subtitle"):
            _text(slide, MARGIN, Inches(4.6), left_w, Inches(1.5),
                  data["subtitle"], size=20, color=t.text_muted,
                  line_spacing=1.3)

        # Accent bar
        _rect(slide, MARGIN, Inches(4.4), Inches(1.2), Inches(0.08), t.primary)

        # Footer with subject + grade
        subject = context.get("subject", "")
        grade = context.get("grade_level", "")
        footer = f"{subject}  \u00b7  {grade} {self.l['class_suffix']}" if grade else subject
        if footer:
            _text(slide, MARGIN, SLIDE_H - Inches(0.5), left_w, Inches(0.3),
                  footer, size=11, color=t.text_muted)

        _speaker_notes(slide, f"{self.l['notes_intro']}: {data.get('title', '')}")
        return slide

    # ── OBJECTIVES ────────────────────────────────────────────────────────
    def build_objectives(self, data: dict, page_num: int, total: int, deck_title: str):
        slide = self._blank()
        t = self.t
        self._motif(slide)

        _text(slide, MARGIN, Inches(0.7), SLIDE_W - 2 * MARGIN, Inches(0.4),
              "01", size=14, color=t.primary, bold=True)
        _text(slide, MARGIN, Inches(1.05), SLIDE_W - 2 * MARGIN, Inches(0.9),
              data.get("title", self.l["objectives_title"]),
              size=38, color=t.text, bold=True)

        items = data.get("items", [])[:4]
        row_h = Inches(0.95)
        start_y = Inches(2.6)

        for i, item in enumerate(items):
            y = start_y + i * (row_h + Inches(0.15))
            circle_d = Inches(0.7)
            cx = MARGIN + circle_d / 2
            cy = y + circle_d / 2
            _circle(slide, cx, cy, circle_d, t.primary)
            _badge_text(slide, cx, cy, str(i + 1),
                        size=20, color=t.text_on_primary)
            _text(slide, MARGIN + Inches(1.0), y + Inches(0.1),
                  SLIDE_W - 2 * MARGIN - Inches(1.0), Inches(0.8),
                  item, size=18, color=t.text,
                  anchor=MSO_ANCHOR.MIDDLE, line_spacing=1.3)

        self._footer(slide, page_num, total, deck_title)
        _speaker_notes(slide, self.l["notes_objectives"])
        return slide

    # ── CONTENT ───────────────────────────────────────────────────────────
    def build_content(self, data: dict, page_num: int, total: int, deck_title: str):
        layout = data.get("layout_hint", "image_left")
        if layout == "stat_callout" and data.get("stat_value"):
            return self._content_stat(data, page_num, total, deck_title)
        elif layout == "image_right":
            return self._content_image(data, page_num, total, deck_title, side="right")
        else:
            return self._content_image(data, page_num, total, deck_title, side="left")

    def _content_image(self, data, page_num, total, deck_title, side="left"):
        slide = self._blank()
        t = self.t
        self._motif(slide)

        img_w = Inches(5.6)
        img_h = Inches(5.6)
        img_y = Inches(1.0)

        if side == "left":
            img_x = MARGIN
            text_x = img_x + img_w + GUTTER
            text_w = SLIDE_W - text_x - MARGIN
        else:
            img_x = SLIDE_W - MARGIN - img_w
            text_x = MARGIN
            text_w = img_x - text_x - GUTTER

        # Image
        img_bytes = self.images.fetch(data.get("image_query", ""))
        if img_bytes:
            _image_fill(slide, img_bytes, img_x, img_y, img_w, img_h)
        else:
            _image_placeholder(slide, img_x, img_y, img_w, img_h, t)

        # Eyebrow
        _text(slide, text_x, Inches(1.0), text_w, Inches(0.4),
              self.l["topic"], size=12, color=t.primary, bold=True)
        # Title
        _text(slide, text_x, Inches(1.4), text_w, Inches(2.0),
              data.get("title", ""), size=30, color=t.text, bold=True,
              line_spacing=1.1)
        # Accent bar
        _rect(slide, text_x, Inches(3.5), Inches(1.0), Inches(0.06), t.primary)
        # Body
        _text(slide, text_x, Inches(3.8), text_w, Inches(3.0),
              data.get("body", ""), size=16, color=t.text, line_spacing=1.5)

        self._footer(slide, page_num, total, deck_title)
        _speaker_notes(slide, f"{self.l['notes_discuss']}: {data.get('title', '')}")
        return slide

    def _content_stat(self, data, page_num, total, deck_title):
        """Big number + label + supporting text."""
        slide = self._blank()
        t = self.t
        self._motif(slide)

        # Title
        _text(slide, MARGIN, Inches(0.8), SLIDE_W - 2 * MARGIN, Inches(0.4),
              self.l["key_fact"], size=12, color=t.primary, bold=True)
        _text(slide, MARGIN, Inches(1.2), SLIDE_W - 2 * MARGIN, Inches(0.9),
              data.get("title", ""), size=30, color=t.text, bold=True)

        # Stat card — left half
        stat_x = MARGIN
        stat_y = Inches(2.6)
        stat_w = Inches(5.6)
        stat_h = Inches(4.0)
        _rect(slide, stat_x, stat_y, stat_w, stat_h, t.primary,
              rounded=True, corner_radius=0.05)

        # Auto-scale stat value font size
        stat_value = str(data.get("stat_value", ""))
        n = len(stat_value)
        if n <= 4:
            stat_size = 110
        elif n <= 5:
            stat_size = 72
        elif n <= 7:
            stat_size = 54
        else:
            stat_size = 40

        # Number in upper ~65% of card
        _text(slide, stat_x + Inches(0.25), stat_y + Inches(0.3),
              stat_w - Inches(0.5), Inches(2.6),
              stat_value, size=stat_size, color=t.text_on_primary, bold=True,
              align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE, line_spacing=1.0)
        # Label below
        _text(slide, stat_x + Inches(0.4), stat_y + Inches(3.05),
              stat_w - Inches(0.8), Inches(0.85),
              data.get("stat_label", ""), size=18, color=t.text_on_primary,
              align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.TOP, line_spacing=1.25)

        # Body — right half
        body_x = stat_x + stat_w + GUTTER
        body_w = SLIDE_W - body_x - MARGIN
        _text(slide, body_x, stat_y + Inches(0.3), body_w, Inches(3.6),
              data.get("body", ""), size=16, color=t.text, line_spacing=1.55)

        self._footer(slide, page_num, total, deck_title)
        _speaker_notes(slide, f"{self.l['notes_fact']}: {data.get('stat_value')} — "
                              f"{data.get('stat_label')}")
        return slide

    # ── KEY TERMS ─────────────────────────────────────────────────────────
    def build_key_terms(self, data: dict, page_num: int, total: int, deck_title: str):
        """Card grid — each term-definition in its own card."""
        slide = self._blank()
        t = self.t
        self._motif(slide)

        _text(slide, MARGIN, Inches(0.7), SLIDE_W - 2 * MARGIN, Inches(0.4),
              self.l["terms"], size=12, color=t.primary, bold=True)
        _text(slide, MARGIN, Inches(1.05), SLIDE_W - 2 * MARGIN, Inches(0.9),
              data.get("title", self.l["terms_title"]),
              size=32, color=t.text, bold=True)

        terms = data.get("terms", [])[:6]
        cols = 3
        rows = (len(terms) + cols - 1) // cols

        grid_x = MARGIN
        grid_y = Inches(2.4)
        grid_w = SLIDE_W - 2 * MARGIN
        grid_h = SLIDE_H - grid_y - Inches(0.7)

        gap = Inches(0.25)
        card_w = (grid_w - gap * (cols - 1)) / cols
        card_h = (grid_h - gap * (max(rows, 1) - 1)) / max(rows, 1)

        for i, term_obj in enumerate(terms):
            row = i // cols
            col = i % cols
            x = grid_x + col * (card_w + gap)
            y = grid_y + row * (card_h + gap)

            _rect(slide, x, y, card_w, card_h, t.surface,
                  rounded=True, corner_radius=0.06)
            # Top accent strip
            _rect(slide, x, y, card_w, Inches(0.12), t.primary,
                  rounded=True, corner_radius=0.5)
            # Term
            term = term_obj.get("term", "") if isinstance(term_obj, dict) else str(term_obj)
            defn = term_obj.get("definition", "") if isinstance(term_obj, dict) else ""
            _text(slide, x + Inches(0.25), y + Inches(0.3),
                  card_w - Inches(0.5), Inches(0.6),
                  term, size=16, color=t.primary, bold=True, line_spacing=1.1)
            _text(slide, x + Inches(0.25), y + Inches(0.95),
                  card_w - Inches(0.5), card_h - Inches(1.1),
                  defn, size=12, color=t.text, line_spacing=1.35)

        self._footer(slide, page_num, total, deck_title)
        _speaker_notes(slide, self.l["notes_terms"])
        return slide

    # ── QUIZ ──────────────────────────────────────────────────────────────
    def build_quiz(self, data: dict, page_num: int, total: int, deck_title: str):
        slide = self._blank()
        t = self.t
        self._motif(slide)

        _text(slide, MARGIN, Inches(0.7), SLIDE_W - 2 * MARGIN, Inches(0.4),
              self.l["quiz"], size=12, color=t.primary, bold=True)
        _text(slide, MARGIN, Inches(1.1), SLIDE_W - 2 * MARGIN, Inches(2.0),
              data.get("question", ""), size=28, color=t.text, bold=True,
              line_spacing=1.2)

        options = data.get("options", [])[:4]
        letters = ["A", "B", "C", "D"]

        grid_y = Inches(2.9)
        grid_w = SLIDE_W - 2 * MARGIN
        gap = Inches(0.3)
        card_w = (grid_w - gap) / 2
        card_h = Inches(1.5)

        for i, opt in enumerate(options):
            row = i // 2
            col = i % 2
            x = MARGIN + col * (card_w + gap)
            y = grid_y + row * (card_h + gap)

            # All cards look identical — answer in speaker notes only
            _rect(slide, x, y, card_w, card_h, t.surface,
                  rounded=True, corner_radius=0.08)

            badge_d = Inches(0.85)
            badge_x = x + Inches(0.35)
            badge_y = y + (card_h - badge_d) / 2
            cx = badge_x + badge_d / 2
            cy = badge_y + badge_d / 2
            _circle(slide, cx, cy, badge_d, t.primary)
            _badge_text(slide, cx, cy, letters[i],
                        size=22, color=t.text_on_primary)

            _text(slide, badge_x + badge_d + Inches(0.3), y,
                  card_w - badge_d - Inches(0.8), card_h,
                  opt, size=16, color=t.text,
                  anchor=MSO_ANCHOR.MIDDLE, line_spacing=1.3)

        self._footer(slide, page_num, total, deck_title)
        answer_idx = data.get("answer", 0)
        if isinstance(answer_idx, int) and 0 <= answer_idx < len(options):
            _speaker_notes(
                slide,
                f"{self.l['notes_quiz_answer']}: {letters[answer_idx]} — {options[answer_idx]}\n"
                f"{self.l['notes_quiz']}"
            )
        return slide

    # ── SUMMARY ───────────────────────────────────────────────────────────
    def build_summary(self, data: dict, page_num: int, total: int, deck_title: str):
        slide = self._blank()
        t = self.t

        # Dark sandwich — primary background
        _rect(slide, 0, 0, SLIDE_W, SLIDE_H, t.primary)

        _text(slide, MARGIN, Inches(0.9), SLIDE_W - 2 * MARGIN, Inches(0.5),
              self.l["summary"], size=14, color=t.text_on_primary, bold=True)
        _text(slide, MARGIN, Inches(1.4), SLIDE_W - 2 * MARGIN, Inches(1.1),
              data.get("title", self.l["summary_title"]),
              size=44, color=t.text_on_primary, bold=True)

        items = data.get("items", [])[:5]
        start_y = Inches(3.0)
        item_h = Inches(0.7)
        gap = Inches(0.18)

        for i, item in enumerate(items):
            y = start_y + i * (item_h + gap)
            _rect(slide, MARGIN, y, SLIDE_W - 2 * MARGIN, item_h,
                  t.bg, rounded=True, corner_radius=0.3)
            badge = Inches(0.45)
            _circle(slide, MARGIN + Inches(0.45), y + item_h / 2, badge, t.accent)
            _badge_text(slide, MARGIN + Inches(0.45), y + item_h / 2,
                        "\u2713", size=16, color=t.text)
            _text(slide, MARGIN + Inches(1.05), y,
                  SLIDE_W - 2 * MARGIN - Inches(1.3), item_h,
                  item, size=16, color=t.text,
                  anchor=MSO_ANCHOR.MIDDLE, line_spacing=1.3)

        _speaker_notes(slide, self.l["notes_summary"])
        return slide

    # ── FALLBACK ──────────────────────────────────────────────────────────
    def build_fallback(self, data: dict, page_num: int, total: int, deck_title: str):
        slide = self._blank()
        self._motif(slide)
        _text(slide, MARGIN, Inches(1.0), SLIDE_W - 2 * MARGIN, Inches(1.0),
              data.get("title", ""), size=30, color=self.t.text, bold=True)
        body = data.get("body", "") or data.get("question", "") or str(data.get("items", ""))
        _text(slide, MARGIN, Inches(2.5), SLIDE_W - 2 * MARGIN, Inches(4.0),
              body, size=16, color=self.t.text, line_spacing=1.5)
        self._footer(slide, page_num, total, deck_title)
        return slide


# ─────────────────────────────────────────────────────────────────────────────
# PUBLIC API — matches v1 signature for drop-in switch
# ─────────────────────────────────────────────────────────────────────────────


def _resolve_theme(slides_data: dict, context_data: dict, template: str | None) -> Theme:
    """Pick theme from context_data.theme > template param > default."""
    name = (context_data or {}).get("theme") or template or "warm"
    name = THEME_ALIASES.get(name, name)
    return THEMES.get(name, THEMES["warm"])


def get_available_templates() -> list[dict]:
    """Return v2 themes for frontend (if needed)."""
    return [
        {"slug": name, "label_ru": t.label_ru, "label_kk": t.label_kk}
        for name, t in THEMES.items()
    ]


def export_to_pptx(
    slides_data: dict,
    context_data: dict,
    template: str | None = None,
) -> io.BytesIO:
    """
    Build PPTX from slides_data + context_data.
    Drop-in replacement for v1 export_to_pptx — same signature, same return type.
    """
    theme = _resolve_theme(slides_data, context_data, template)
    l10n = _get_l10n(context_data)
    images = ImageProvider()

    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H

    builder = SlideBuilder(prs, theme, images, l10n)
    deck_title = slides_data.get("title", "") or (context_data or {}).get("paragraph_title", "")
    slides = slides_data.get("slides", [])
    total = len(slides)

    for idx, sd in enumerate(slides):
        page_num = idx + 1
        stype = sd.get("type", "content")

        try:
            if stype == "title":
                builder.build_title(sd, context_data or {}, deck_title)
            elif stype == "objectives":
                builder.build_objectives(sd, page_num, total, deck_title)
            elif stype == "content":
                builder.build_content(sd, page_num, total, deck_title)
            elif stype == "key_terms":
                builder.build_key_terms(sd, page_num, total, deck_title)
            elif stype == "quiz":
                builder.build_quiz(sd, page_num, total, deck_title)
            elif stype == "summary":
                builder.build_summary(sd, page_num, total, deck_title)
            else:
                builder.build_content(sd, page_num, total, deck_title)
        except Exception:
            logger.exception("Error building slide %d type=%s", idx, stype)
            builder.build_fallback(sd, page_num, total, deck_title)

    buf = io.BytesIO()
    prs.save(buf)
    buf.seek(0)
    return buf
