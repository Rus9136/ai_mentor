"""
PPTX export with built-in design system.

Blank-canvas approach: every slide is built from shapes with a consistent
color palette, typography, spacing, and visual elements (accent bars,
card backgrounds, decorative shapes).

No template placeholders — full control over every pixel.
"""
import logging
from io import BytesIO
from pathlib import Path
from typing import Optional

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt

from app.core.config import settings

logger = logging.getLogger(__name__)

# ═══════════════════════════════════════════════════════════
#  DESIGN SYSTEM
# ═══════════════════════════════════════════════════════════

# Slide canvas (16:9 widescreen)
SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)

# Margins & grid
MARGIN = Inches(0.9)
CONTENT_L = MARGIN
CONTENT_T = Inches(1.7)          # Below header
CONTENT_W = SLIDE_W - MARGIN * 2
CONTENT_H = SLIDE_H - CONTENT_T - MARGIN

# Color themes
THEMES = {
    "blue": {
        "primary":    RGBColor(0x1E, 0x3A, 0x5F),   # Deep navy
        "secondary":  RGBColor(0x2D, 0x6A, 0x9F),   # Medium blue
        "accent":     RGBColor(0x3B, 0x82, 0xF6),   # Bright blue
        "accent2":    RGBColor(0x10, 0xB9, 0x81),   # Teal green
        "bg_light":   RGBColor(0xF0, 0xF4, 0xFA),   # Very light blue
        "bg_card":    RGBColor(0xE8, 0xEE, 0xF6),   # Card background
        "bg_correct": RGBColor(0xD1, 0xFA, 0xE5),   # Green for correct
        "text_dark":  RGBColor(0x1E, 0x29, 0x3B),   # Near-black
        "text_mid":   RGBColor(0x64, 0x74, 0x8B),   # Gray
        "text_white": RGBColor(0xFF, 0xFF, 0xFF),   # White
        "divider":    RGBColor(0xCB, 0xD5, 0xE1),   # Light gray line
    },
    "green": {
        "primary":    RGBColor(0x14, 0x53, 0x2D),
        "secondary":  RGBColor(0x16, 0x6B, 0x3A),
        "accent":     RGBColor(0x22, 0xC5, 0x5E),
        "accent2":    RGBColor(0x0E, 0xA5, 0xE9),
        "bg_light":   RGBColor(0xF0, 0xFD, 0xF4),
        "bg_card":    RGBColor(0xDC, 0xFC, 0xE7),
        "bg_correct": RGBColor(0xBB, 0xF7, 0xD0),
        "text_dark":  RGBColor(0x1E, 0x29, 0x3B),
        "text_mid":   RGBColor(0x64, 0x74, 0x8B),
        "text_white": RGBColor(0xFF, 0xFF, 0xFF),
        "divider":    RGBColor(0xBB, 0xF7, 0xD0),
    },
    "warm": {
        "primary":    RGBColor(0x7C, 0x2D, 0x12),
        "secondary":  RGBColor(0xC2, 0x41, 0x0C),
        "accent":     RGBColor(0xF9, 0x73, 0x16),
        "accent2":    RGBColor(0x06, 0x84, 0x5D),
        "bg_light":   RGBColor(0xFF, 0xF7, 0xED),
        "bg_card":    RGBColor(0xFF, 0xED, 0xD5),
        "bg_correct": RGBColor(0xD1, 0xFA, 0xE5),
        "text_dark":  RGBColor(0x1E, 0x29, 0x3B),
        "text_mid":   RGBColor(0x78, 0x71, 0x6C),
        "text_white": RGBColor(0xFF, 0xFF, 0xFF),
        "divider":    RGBColor(0xFE, 0xD7, 0xAA),
    },
}

# Typography
FONT_FAMILY = "Calibri"
FONT_TITLE = Pt(36)
FONT_HEADING = Pt(28)
FONT_SUBHEADING = Pt(20)
FONT_BODY = Pt(17)
FONT_SMALL = Pt(14)
FONT_CAPTION = Pt(12)

# Header bar
HEADER_H = Inches(1.3)

# Card styling
CARD_RADIUS = Inches(0.12)
CARD_PAD = Inches(0.25)


def get_available_templates() -> list[dict]:
    """Return color themes as 'templates' for the frontend."""
    return [
        {"slug": "blue", "label": "Академический (синий)"},
        {"slug": "green", "label": "Биология (зелёный)"},
        {"slug": "warm", "label": "История (тёплый)"},
    ]


def export_to_pptx(slides_data: dict, context_data: dict, template: str | None = None) -> BytesIO:
    """Build a presentation from scratch with the design system."""
    theme_name = template if template in THEMES else "blue"
    T = THEMES[theme_name]
    textbook_id = context_data.get("textbook_id")

    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H
    blank = prs.slide_layouts[6]  # Blank layout

    slides = slides_data.get("slides", [])
    total = len(slides)

    for idx, sd in enumerate(slides):
        slide = prs.slides.add_slide(blank)
        st = sd.get("type", "content")

        try:
            if st == "title":
                _build_title(slide, sd, context_data, T)
            elif st == "objectives":
                _build_bullets(slide, sd, T, numbered=True)
            elif st == "content":
                _build_content(slide, sd, T, textbook_id)
            elif st == "key_terms":
                _build_key_terms(slide, sd, T)
            elif st == "quiz":
                _build_quiz(slide, sd, T)
            elif st == "summary":
                _build_summary(slide, sd, T)
            else:
                _build_content(slide, sd, T, textbook_id)
        except Exception:
            logger.exception("Error building slide %d type=%s", idx, st)
            _build_fallback(slide, sd, T)

        # Slide number (skip title slide)
        if st != "title":
            _add_slide_number(slide, idx + 1, total, T)

    buf = BytesIO()
    prs.save(buf)
    buf.seek(0)
    return buf


# ═══════════════════════════════════════════════════════════
#  SLIDE BUILDERS
# ═══════════════════════════════════════════════════════════


def _build_title(slide, data, context, T):
    """Title slide: full-color background, large centered text."""
    # Full background
    _rect(slide, 0, 0, SLIDE_W, SLIDE_H, T["primary"])

    # Decorative accent bar at top
    _rect(slide, 0, 0, SLIDE_W, Inches(0.08), T["accent"])

    # Title
    title = data.get("title", "")
    _text(slide, Inches(1.2), Inches(1.8), SLIDE_W - Inches(2.4), Inches(2.5),
          title, FONT_TITLE, T["text_white"], bold=True, align=PP_ALIGN.CENTER)

    # Subtitle
    subtitle = data.get("subtitle", "")
    if subtitle:
        _text(slide, Inches(1.5), Inches(4.3), SLIDE_W - Inches(3.0), Inches(1.0),
              subtitle, FONT_SUBHEADING, RGBColor(0xBF, 0xDB, 0xFE), align=PP_ALIGN.CENTER)

    # Bottom bar with context
    _rect(slide, 0, SLIDE_H - Inches(0.9), SLIDE_W, Inches(0.9), T["secondary"])
    subject = context.get("subject", "")
    grade = context.get("grade_level", "")
    footer = f"{subject}  |  {grade} класс" if grade else subject
    _text(slide, MARGIN, SLIDE_H - Inches(0.75), CONTENT_W, Inches(0.5),
          footer, FONT_SMALL, RGBColor(0xE2, 0xE8, 0xF0), align=PP_ALIGN.CENTER)

    # Decorative circle
    _circle(slide, SLIDE_W - Inches(2.0), Inches(1.0), Inches(0.6), T["accent"], alpha=50000)


def _build_bullets(slide, data, T, numbered=False):
    """Objectives / bullet list slide with accent sidebar."""
    # Background
    _rect(slide, 0, 0, SLIDE_W, SLIDE_H, T["bg_light"])

    # Left accent sidebar
    _rect(slide, 0, 0, Inches(0.15), SLIDE_H, T["accent"])

    # Header area
    _rect(slide, Inches(0.15), 0, SLIDE_W - Inches(0.15), HEADER_H, T["primary"])
    _text(slide, MARGIN + Inches(0.3), Inches(0.3), CONTENT_W, Inches(0.7),
          data.get("title", ""), FONT_HEADING, T["text_white"], bold=True)

    # Bullet items as cards
    items = data.get("items", [])[:8]
    y = CONTENT_T + Inches(0.2)
    card_h = Inches(0.65)
    gap = Inches(0.12)

    for i, item in enumerate(items):
        # Card background (alternating subtle colors)
        bg = T["bg_card"] if i % 2 == 0 else T["bg_light"]
        _rounded_rect(slide, MARGIN + Inches(0.3), y, CONTENT_W - Inches(0.3), card_h, bg)

        # Number badge or bullet
        badge_w = Inches(0.45)
        if numbered:
            _rounded_rect(slide, MARGIN + Inches(0.45), y + Inches(0.1),
                          badge_w, Inches(0.45), T["accent"])
            _text(slide, MARGIN + Inches(0.45), y + Inches(0.1),
                  badge_w, Inches(0.45), str(i + 1),
                  FONT_BODY, T["text_white"], bold=True, align=PP_ALIGN.CENTER)
        else:
            _text(slide, MARGIN + Inches(0.5), y + Inches(0.08),
                  Inches(0.3), card_h, "\u2022",
                  FONT_HEADING, T["accent"], bold=True)

        # Item text
        text_left = MARGIN + Inches(1.1)
        _text(slide, text_left, y + Inches(0.12),
              CONTENT_W - Inches(1.2), Inches(0.42),
              item, FONT_BODY, T["text_dark"])

        y += card_h + gap


def _build_content(slide, data, T, textbook_id):
    """Content slide with header + body text + optional image."""
    # Background
    _rect(slide, 0, 0, SLIDE_W, SLIDE_H, T["bg_light"])

    # Header bar
    _rect(slide, 0, 0, SLIDE_W, HEADER_H, T["primary"])
    _text(slide, MARGIN, Inches(0.3), CONTENT_W, Inches(0.7),
          data.get("title", ""), FONT_HEADING, T["text_white"], bold=True)

    # Accent line under header
    _rect(slide, MARGIN, HEADER_H + Inches(0.05), Inches(2.5), Inches(0.05), T["accent"])

    # Check for image
    image_url = data.get("image_url")
    image_path = _resolve_image(image_url, textbook_id) if image_url and textbook_id else None
    has_image = image_path and Path(image_path).exists()

    if has_image:
        # Two-column: text left, image right
        text_w = CONTENT_W * 0.55
        _text(slide, MARGIN, CONTENT_T + Inches(0.15), text_w, CONTENT_H - Inches(0.3),
              data.get("body", ""), FONT_BODY, T["text_dark"], wrap=True, line_spacing=1.3)

        # Image card
        img_l = MARGIN + text_w + Inches(0.4)
        img_w = CONTENT_W - text_w - Inches(0.4)
        img_h = min(img_w * 0.75, CONTENT_H - Inches(0.3))
        # Card shadow background
        _rounded_rect(slide, img_l - Inches(0.05), CONTENT_T + Inches(0.1),
                      img_w + Inches(0.1), img_h + Inches(0.1), T["bg_card"])
        try:
            slide.shapes.add_picture(image_path, img_l, CONTENT_T + Inches(0.15),
                                     img_w, img_h)
        except Exception:
            logger.warning("Failed to add image: %s", image_path)
    else:
        # Full-width content card
        _rounded_rect(slide, MARGIN, CONTENT_T + Inches(0.15),
                      CONTENT_W, CONTENT_H - Inches(0.3), RGBColor(0xFF, 0xFF, 0xFF))
        _text(slide, MARGIN + CARD_PAD, CONTENT_T + Inches(0.35),
              CONTENT_W - CARD_PAD * 2, CONTENT_H - Inches(0.7),
              data.get("body", ""), FONT_BODY, T["text_dark"], wrap=True, line_spacing=1.3)


def _build_key_terms(slide, data, T):
    """Key terms as styled cards in a grid."""
    # Background
    _rect(slide, 0, 0, SLIDE_W, SLIDE_H, T["bg_light"])

    # Header
    _rect(slide, 0, 0, SLIDE_W, HEADER_H, T["primary"])
    _text(slide, MARGIN, Inches(0.3), CONTENT_W, Inches(0.7),
          data.get("title", ""), FONT_HEADING, T["text_white"], bold=True)

    terms = data.get("terms", [])[:6]
    y = CONTENT_T + Inches(0.15)
    card_h = Inches(0.85)
    gap = Inches(0.12)
    term_w = Inches(3.2)

    for i, term_data in enumerate(terms):
        term = term_data.get("term", "")
        defn = term_data.get("definition", "")

        # Full-width card
        _rounded_rect(slide, MARGIN, y, CONTENT_W, card_h, RGBColor(0xFF, 0xFF, 0xFF))

        # Left accent stripe on card
        _rect(slide, MARGIN, y, Inches(0.08), card_h, T["accent"])

        # Term (bold, colored)
        _text(slide, MARGIN + Inches(0.25), y + Inches(0.15),
              term_w, Inches(0.55), term, FONT_BODY, T["secondary"], bold=True)

        # Separator dot
        _text(slide, MARGIN + term_w + Inches(0.1), y + Inches(0.15),
              Inches(0.3), Inches(0.55), "\u2014", FONT_BODY, T["divider"])

        # Definition
        _text(slide, MARGIN + term_w + Inches(0.5), y + Inches(0.15),
              CONTENT_W - term_w - Inches(0.8), Inches(0.55),
              defn, FONT_BODY, T["text_dark"])

        y += card_h + gap


def _build_quiz(slide, data, T):
    """Quiz slide with question card + styled option cards."""
    # Background
    _rect(slide, 0, 0, SLIDE_W, SLIDE_H, T["bg_light"])

    # Header with accent2 color
    _rect(slide, 0, 0, SLIDE_W, HEADER_H, T["accent2"])
    _text(slide, MARGIN, Inches(0.3), CONTENT_W, Inches(0.7),
          data.get("title", "\u2753 Білімді тексер"), FONT_HEADING, T["text_white"], bold=True)

    # Question card
    question = data.get("question", "")
    q_h = Inches(1.1)
    _rounded_rect(slide, MARGIN, CONTENT_T, CONTENT_W, q_h, RGBColor(0xFF, 0xFF, 0xFF))
    _rect(slide, MARGIN, CONTENT_T, Inches(0.08), q_h, T["accent2"])
    _text(slide, MARGIN + Inches(0.3), CONTENT_T + Inches(0.15),
          CONTENT_W - Inches(0.6), Inches(0.8),
          question, FONT_SUBHEADING, T["text_dark"], bold=True, wrap=True)

    # Option cards
    options = data.get("options", [])[:6]
    answer_idx = data.get("answer", 0)
    letters = ["A", "B", "C", "D", "E", "F"]

    y = CONTENT_T + q_h + Inches(0.25)
    opt_h = Inches(0.6)
    gap = Inches(0.1)

    for i, option in enumerate(options):
        is_correct = (i == answer_idx)

        # Option card
        bg = T["bg_correct"] if is_correct else RGBColor(0xFF, 0xFF, 0xFF)
        _rounded_rect(slide, MARGIN + Inches(0.2), y, CONTENT_W - Inches(0.4), opt_h, bg)

        # Letter badge
        badge_color = T["accent2"] if is_correct else T["text_mid"]
        _rounded_rect(slide, MARGIN + Inches(0.35), y + Inches(0.08),
                      Inches(0.44), Inches(0.44), badge_color)
        _text(slide, MARGIN + Inches(0.35), y + Inches(0.08),
              Inches(0.44), Inches(0.44), letters[i],
              FONT_BODY, T["text_white"], bold=True, align=PP_ALIGN.CENTER)

        # Option text
        text_color = T["text_dark"]
        _text(slide, MARGIN + Inches(1.0), y + Inches(0.1),
              CONTENT_W - Inches(1.6), Inches(0.4),
              option, FONT_BODY, text_color, bold=is_correct)

        # Checkmark for correct answer
        if is_correct:
            _text(slide, CONTENT_W + MARGIN - Inches(0.8), y + Inches(0.08),
                  Inches(0.5), Inches(0.44), "\u2714",
                  FONT_SUBHEADING, T["accent2"], bold=True)

        y += opt_h + gap


def _build_summary(slide, data, T):
    """Summary slide with checklist items and accent styling."""
    # Background gradient effect: primary top, light bottom
    _rect(slide, 0, 0, SLIDE_W, Inches(2.5), T["primary"])
    _rect(slide, 0, Inches(2.5), SLIDE_W, SLIDE_H - Inches(2.5), T["bg_light"])

    # Title in top area
    _text(slide, MARGIN, Inches(0.5), CONTENT_W, Inches(1.0),
          data.get("title", "Қорытынды"), FONT_HEADING, T["text_white"], bold=True)

    # Subtitle line
    _text(slide, MARGIN, Inches(1.3), CONTENT_W, Inches(0.5),
          "Негізгі тұжырымдар", FONT_SMALL, RGBColor(0xBF, 0xDB, 0xFE))

    # Summary cards
    items = data.get("items", [])[:8]
    y = Inches(2.8)
    card_h = Inches(0.6)
    gap = Inches(0.1)

    for i, item in enumerate(items):
        # Card
        _rounded_rect(slide, MARGIN, y, CONTENT_W, card_h, RGBColor(0xFF, 0xFF, 0xFF))

        # Checkmark circle
        _circle(slide, MARGIN + Inches(0.15), y + Inches(0.1),
                Inches(0.4), T["accent2"])
        _text(slide, MARGIN + Inches(0.15), y + Inches(0.1),
              Inches(0.4), Inches(0.4), "\u2713",
              FONT_SMALL, T["text_white"], bold=True, align=PP_ALIGN.CENTER)

        # Item text
        _text(slide, MARGIN + Inches(0.75), y + Inches(0.12),
              CONTENT_W - Inches(1.0), Inches(0.38),
              item, FONT_BODY, T["text_dark"])

        y += card_h + gap


def _build_fallback(slide, data, T):
    """Simple fallback slide."""
    _rect(slide, 0, 0, SLIDE_W, HEADER_H, T["primary"])
    _text(slide, MARGIN, Inches(0.3), CONTENT_W, Inches(0.7),
          data.get("title", ""), FONT_HEADING, T["text_white"], bold=True)
    body = data.get("body", "") or data.get("question", "") or str(data.get("items", ""))
    _text(slide, MARGIN, CONTENT_T, CONTENT_W, CONTENT_H,
          body, FONT_BODY, T["text_dark"], wrap=True)


# ═══════════════════════════════════════════════════════════
#  SHAPE PRIMITIVES
# ═══════════════════════════════════════════════════════════


def _rect(slide, left, top, width, height, fill_color):
    """Solid rectangle (no border)."""
    shape = slide.shapes.add_shape(1, int(left), int(top), int(width), int(height))
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    shape.line.fill.background()
    return shape


def _rounded_rect(slide, left, top, width, height, fill_color):
    """Rounded rectangle."""
    from pptx.enum.shapes import MSO_SHAPE
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, int(left), int(top), int(width), int(height)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    shape.line.fill.background()
    return shape


def _circle(slide, left, top, diameter, fill_color, alpha=None):
    """Circle shape."""
    from pptx.enum.shapes import MSO_SHAPE
    shape = slide.shapes.add_shape(
        MSO_SHAPE.OVAL, int(left), int(top), int(diameter), int(diameter)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    shape.line.fill.background()
    return shape


def _text(slide, left, top, width, height, text, font_size, font_color,
          bold=False, align=PP_ALIGN.LEFT, wrap=True, line_spacing=None):
    """Text box with consistent styling."""
    txBox = slide.shapes.add_textbox(int(left), int(top), int(width), int(height))
    tf = txBox.text_frame
    tf.word_wrap = wrap
    tf.auto_size = None
    tf.margin_left = Pt(2)
    tf.margin_right = Pt(2)
    tf.margin_top = Pt(1)
    tf.margin_bottom = Pt(1)

    p = tf.paragraphs[0]
    p.alignment = align

    if line_spacing:
        p.line_spacing = line_spacing

    run = p.add_run()
    run.text = text
    run.font.name = FONT_FAMILY
    run.font.size = font_size
    run.font.color.rgb = font_color
    run.font.bold = bold

    return txBox


def _add_slide_number(slide, num, total, T):
    """Small slide number in bottom-right corner."""
    _text(slide, SLIDE_W - Inches(1.0), SLIDE_H - Inches(0.45),
          Inches(0.7), Inches(0.3),
          f"{num}/{total}", FONT_CAPTION, T["text_mid"], align=PP_ALIGN.RIGHT)


# ═══════════════════════════════════════════════════════════
#  UTILS
# ═══════════════════════════════════════════════════════════


def _resolve_image(url: str, textbook_id: int) -> Optional[str]:
    """Convert /uploads/textbook-images/ID/file.jpg to absolute path."""
    if not url or not url.startswith("/uploads/textbook-images/"):
        return None
    filename = url.split("/")[-1]
    path = Path(settings.UPLOAD_DIR) / "textbook-images" / str(textbook_id) / filename
    return str(path) if path.exists() else None
